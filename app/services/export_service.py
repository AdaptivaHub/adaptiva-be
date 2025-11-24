from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from fastapi import HTTPException
from fastapi.responses import FileResponse
import io
import json
from pathlib import Path
from datetime import datetime

from app.utils import get_dataframe
from app.models import ExportRequest, ExportFormat
from app.services.insights_service import get_data_insights


def export_to_pdf(request: ExportRequest) -> str:
    """
    Export data to PDF format
    
    Args:
        request: ExportRequest with export parameters
        
    Returns:
        Path to generated PDF file
    """
    try:
        # Get the dataframe
        df = get_dataframe(request.file_id)
        
        # Create PDF file
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"export_{request.file_id[:8]}_{timestamp}.pdf"
        filepath = Path("uploads") / filename
        
        # Create PDF document
        doc = SimpleDocTemplate(str(filepath), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        title = Paragraph("Data Analysis Report", title_style)
        story.append(title)
        story.append(Spacer(1, 0.2 * inch))
        
        # Dataset info
        info_style = styles['Normal']
        story.append(Paragraph(f"<b>Dataset Information</b>", styles['Heading2']))
        story.append(Spacer(1, 0.1 * inch))
        story.append(Paragraph(f"Rows: {len(df)}", info_style))
        story.append(Paragraph(f"Columns: {len(df.columns)}", info_style))
        story.append(Spacer(1, 0.2 * inch))
        
        # Add insights if requested
        if request.include_insights:
            insights = get_data_insights(request.file_id)
            
            story.append(Paragraph("<b>Data Insights</b>", styles['Heading2']))
            story.append(Spacer(1, 0.1 * inch))
            
            # Column information table
            col_data = [["Column", "Type", "Non-Null", "Unique Values"]]
            for col, info in insights.column_info.items():
                col_data.append([
                    col,
                    info['dtype'],
                    str(info['non_null_count']),
                    str(info['unique_values'])
                ])
            
            col_table = Table(col_data, repeatRows=1)
            col_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
            ]))
            story.append(col_table)
            story.append(Spacer(1, 0.2 * inch))
            
            # Missing values
            if sum(insights.missing_values.values()) > 0:
                story.append(Paragraph("<b>Missing Values</b>", styles['Heading3']))
                story.append(Spacer(1, 0.1 * inch))
                for col, count in insights.missing_values.items():
                    if count > 0:
                        story.append(Paragraph(f"{col}: {count}", info_style))
                story.append(Spacer(1, 0.2 * inch))
        
        # Data preview
        story.append(Paragraph("<b>Data Preview (First 10 Rows)</b>", styles['Heading2']))
        story.append(Spacer(1, 0.1 * inch))
        
        # Limit to first 10 rows and 6 columns for readability
        preview_df = df.head(10).iloc[:, :6]
        
        # Create table data
        table_data = [preview_df.columns.tolist()]
        for _, row in preview_df.iterrows():
            table_data.append([str(val)[:20] for val in row.values])  # Truncate long values
        
        data_table = Table(table_data, repeatRows=1)
        data_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('FONTSIZE', (0, 1), (-1, -1), 7),
        ]))
        story.append(data_table)
        
        # Build PDF
        doc.build(story)
        
        return str(filepath)
        
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error exporting to PDF: {str(e)}")


def export_to_pptx(request: ExportRequest) -> str:
    """
    Export data to PowerPoint format
    
    Args:
        request: ExportRequest with export parameters
        
    Returns:
        Path to generated PPTX file
    """
    try:
        # Get the dataframe
        df = get_dataframe(request.file_id)
        
        # Create PPTX file
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"export_{request.file_id[:8]}_{timestamp}.pptx"
        filepath = Path("uploads") / filename
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Data Analysis Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Dataset information slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Dataset Information"
        tf = body_shape.text_frame
        tf.text = f"Rows: {len(df)}"
        
        p = tf.add_paragraph()
        p.text = f"Columns: {len(df.columns)}"
        
        p = tf.add_paragraph()
        p.text = f"Column Names: {', '.join(df.columns.tolist()[:10])}"
        if len(df.columns) > 10:
            p2 = tf.add_paragraph()
            p2.text = "... and more"
        
        # Add insights if requested
        if request.include_insights:
            insights = get_data_insights(request.file_id)
            
            # Insights slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = "Data Insights"
            tf = body_shape.text_frame
            tf.text = f"Total Missing Values: {sum(insights.missing_values.values())}"
            
            p = tf.add_paragraph()
            p.text = f"Duplicate Rows: {insights.duplicates_count}"
            
            p = tf.add_paragraph()
            p.text = f"Numerical Columns: {len(insights.numerical_summary)}"
            
        # Data preview slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        title_shape = slide.shapes.title
        title_shape.text = "Data Preview"
        
        # Add table
        preview_df = df.head(5).iloc[:, :5]  # First 5 rows and 5 columns
        
        rows = len(preview_df) + 1
        cols = len(preview_df.columns)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.3) * rows
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column headings
        for col_idx, col_name in enumerate(preview_df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Fill in data
        for row_idx, row in enumerate(preview_df.itertuples(index=False), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)[:30]  # Truncate long values
                cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Save presentation
        prs.save(str(filepath))
        
        return str(filepath)
        
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error exporting to PPTX: {str(e)}")


def export_data(request: ExportRequest) -> str:
    """
    Export data based on format
    
    Args:
        request: ExportRequest with export parameters
        
    Returns:
        Path to generated file
    """
    if request.export_format == ExportFormat.PDF:
        return export_to_pdf(request)
    elif request.export_format == ExportFormat.PPTX:
        return export_to_pptx(request)
    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported export format: {request.export_format}"
        )
